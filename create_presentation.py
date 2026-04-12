from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_content_slide(prs, title, bullet_points):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = bullet_points[0]
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
    return slide


def add_code_slide(prs, title, code_text):
    layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(layout)

    # Title box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].runs[0].font.size = Pt(28)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Code box
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    tf = code_box.text_frame
    tf.word_wrap = True
    tf.text = code_text
    run = tf.paragraphs[0].runs[0]
    run.font.name = "Courier New"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x24, 0x29, 0x2E)

    return slide


def main():
    prs = Presentation()

    # Slide 1: Title
    add_title_slide(
        prs,
        "Python Scripts Overview",
        "A summary of utility scripts in this repository",
    )

    # Slide 2: Repository Overview
    add_content_slide(
        prs,
        "Repository Contents",
        [
            "expand_pc.py  –  Expands numeric ranges into full lists",
            "create_presentation.py  –  Generates this PowerPoint",
            "Written in Python 3",
            "Lightweight, no external dependencies for core scripts",
        ],
    )

    # Slide 3: expand_pc.py explained
    add_content_slide(
        prs,
        "expand_pc.py – How It Works",
        [
            "Accepts a comma-separated string of numeric ranges (e.g. '1000-1920')",
            "Splits each range on '-' and expands to every integer in that range",
            "Joins all expanded numbers into a single comma-separated string",
            "Example input:  \"1000-1920, 2000-2239\"",
            "Example output: \"1000, 1001, 1002, … 2239\"",
        ],
    )

    # Slide 4: Code sample
    add_code_slide(
        prs,
        "expand_pc.py – Code",
        (
            "def expand_ranges(ranges):\n"
            "    expanded_list = []\n"
            "    for r in ranges.split(', '):\n"
            "        start, end = map(int, r.split('-'))\n"
            "        expanded_list.extend(range(start, end + 1))\n"
            "    return expanded_list\n\n"
            "ranges = '1000-1920, 2000-2239, 2555-2574, 2740-2786'\n"
            "expanded = expand_ranges(ranges)\n"
            "print(', '.join(map(str, expanded)))"
        ),
    )

    # Slide 5: Summary
    add_content_slide(
        prs,
        "Summary",
        [
            "Small, focused scripts for specific utility tasks",
            "Easy to extend with additional range sets",
            "python-pptx used to automate presentation creation",
            "Run create_presentation.py to regenerate this file",
        ],
    )

    output_path = "presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")


if __name__ == "__main__":
    main()
